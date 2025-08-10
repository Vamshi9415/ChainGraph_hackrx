import os
import uvicorn
import asyncio
import logging
import torch
import json
import re
from fastapi import FastAPI, HTTPException, Depends, status
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel
from dotenv import load_dotenv
import requests
from io import BytesIO
from typing import List, Optional, Dict, Any, Union, Set, Tuple
import aiofiles
import uuid
from datetime import datetime
from urllib.parse import urlparse, urljoin, urlunparse
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
import zipfile
import base64
import numpy as np

# --- Document Specific Imports ---
import fitz
import docx
from email import message_from_bytes
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup
import openpyxl
from xml.etree import ElementTree as ET

# --- LangChain Imports ---
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter, CharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain.agents import AgentExecutor, create_tool_calling_agent
from langchain.tools import tool
from langchain_community.document_loaders import WebBaseLoader
from langchain.chains.combine_documents import create_stuff_documents_chain

# --- Enhanced Imports ---
try:
    from sentence_transformers import CrossEncoder
    RERANK_AVAILABLE = True
except ImportError:
    RERANK_AVAILABLE = False
    logging.warning("sentence-transformers not available. Re-ranking disabled.")

try:
    # OCR is now handled by pytesseract
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logging.warning("pytesseract not available. OCR features disabled.")

try:
    import langdetect
    LANG_DETECT_AVAILABLE = True
except ImportError:
    LANG_DETECT_AVAILABLE = False
    logging.warning("langdetect not available. Language detection will use fallback mechanisms.")

# --- Configuration ---
load_dotenv()
AUTH_TOKEN = os.getenv("AUTH_TOKEN", "2b55e57dd2584f97b52854b0738dc5608ab353c4fbc8d0409b738b7b21218fbb")
LOG_REQUESTS_PATH, TEMP_FILES_PATH = "./request_logs", "./temp_files"
os.makedirs(LOG_REQUESTS_PATH, exist_ok=True)
os.makedirs(TEMP_FILES_PATH, exist_ok=True)

# Enhanced logging configuration
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('system.log')
    ]
)

# Device detection for optimal model loading
if torch.cuda.is_available():
    device = "cuda"
elif hasattr(torch.backends, 'mps') and torch.backends.mps.is_available() and torch.backends.mps.is_built():
    device = "mps"
else:
    device = "cpu"

logging.info(f"Using device: {device}")

# --- Enhanced Data Classes ---
@dataclass
class ExtractedURL:
    url: str
    context: str
    source_location: str
    confidence: float
    url_type: str

@dataclass
class ExtractedTable:
    content: str
    table_type: str
    location: str
    metadata: Dict[str, Any]

@dataclass
class ExtractedImage:
    image_path: str
    ocr_text: str
    metadata: Dict[str, Any]
    confidence: float

@dataclass
class ProcessedDocument:
    content: str
    metadata: Dict[str, Any]
    tables: List[ExtractedTable]
    images: List[ExtractedImage]
    detected_language: str
    extracted_urls: List[ExtractedURL]

@dataclass
class ChunkMetadata:
    chunk_id: int
    char_count: int
    word_count: int
    has_tables: bool
    has_urls: bool
    importance_score: float
    content_type: str

# --- Global Models ---
embeddings_fast, embeddings_accurate, reranker, llm = None, None, None, None

def initialize_models():
    global embeddings_fast, embeddings_accurate, reranker, llm

    # Fast embedding model for general use
    if embeddings_fast is None:
        try:
            embeddings_fast = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2",
                model_kwargs={"device": device},
                encode_kwargs={"normalize_embeddings": True}
            )
            logging.info("Fast embeddings model loaded successfully")
        except Exception as e:
            logging.error(f"Failed to load fast embeddings model: {e}")
            embeddings_fast = None

    # Accurate embedding model for complex documents
    if embeddings_accurate is None:
        try:
            embeddings_accurate = HuggingFaceEmbeddings(
                model_name="BAAI/bge-small-en-v1.5",
                model_kwargs={"device": device},
                encode_kwargs={"normalize_embeddings": True}
            )
            logging.info("Accurate embeddings model loaded successfully")
        except Exception as e:
            logging.error(f"Failed to load accurate embeddings model: {e}")
            embeddings_accurate = None

    # Re-ranking model
    if RERANK_AVAILABLE and reranker is None:
        try:
            reranker = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2', device=device)
            logging.info("Re-ranking model loaded successfully")
        except Exception as e:
            logging.warning(f"Failed to initialize re-ranker: {e}")
            reranker = None

    # LLM model
    if llm is None:
        try:
            llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash-lite", temperature=0.0, max_tokens=512)
            logging.info("LLM loaded successfully")
        except Exception as e:
            logging.error(f"Failed to initialize LLM: {e}")
            llm = None

    # OCR reader (pytesseract) does not require a reader object to be initialized
    if OCR_AVAILABLE:
        logging.info("Pytesseract OCR is available.")
    else:
        logging.warning("Pytesseract OCR is not available.")


initialize_models()

# --- FastAPI Setup ---
app = FastAPI(title="Enhanced Language-Strict Document-Targeted RAG System", version="13.0")

class QueryRequest(BaseModel):
    documents: str
    questions: List[str]

class QueryResponse(BaseModel):
    answers: List[str]

bearer_scheme = HTTPBearer()

def validate_token(credentials: HTTPAuthorizationCredentials = Depends(bearer_scheme)):
    """Validate authentication token"""
    if credentials.scheme != "Bearer" or credentials.credentials != AUTH_TOKEN:
        raise HTTPException(status_code=401, detail="Invalid authentication token")
    return True

def setup_request_logger(request_id: str):
    """Set up a dedicated logger for each request"""
    logger = logging.getLogger(request_id)
    if not logger.handlers:
        handler = logging.FileHandler(os.path.join(LOG_REQUESTS_PATH, f"{request_id}.log"))
        formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
        handler.setFormatter(formatter)
        logger.addHandler(handler)
        logger.setLevel(logging.INFO)
    return logger

# --- Enhanced Document Type Detection ---
def detect_document_type_strict(doc_url: str) -> str:
    """Strict document type detection based on file extensions"""
    path = urlparse(doc_url).path.lower()

    # Image types for OCR processing only
    if path.endswith('.png'):
        return "png"
    elif path.endswith(('.jpg', '.jpeg')):
        return "jpeg"

    # Document types
    elif path.endswith('.pdf'):
        return "pdf"
    elif path.endswith(('.docx', '.doc')):
        return "docx"
    elif path.endswith(('.pptx', '.ppt')):
        return "pptx"
    elif path.endswith('.txt'):
        return "txt"
    elif path.endswith(('.xlsx', '.xls')):
        return "xlsx"
    elif path.endswith('.html') or path.endswith('.htm'):
        return "html"

    # Web URLs default to HTML
    elif urlparse(doc_url).scheme in ["http", "https"]:
        return "html"

    return "unknown"

# --- Enhanced Language Detection ---
def detect_language_robust(text: str) -> str:
    """Enhanced language detection with better accuracy"""
    if not LANG_DETECT_AVAILABLE or not text.strip():
        return "en"

    try:
        # Clean text for better detection - remove URLs, numbers, special chars
        cleaned_text = re.sub(r'https?://\S+', '', text)
        cleaned_text = re.sub(r'\d+', '', cleaned_text)
        cleaned_text = re.sub(r'[^\w\s]', ' ', cleaned_text)
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()

        # Use a larger sample for better accuracy
        sample_text = cleaned_text[:5000] if len(cleaned_text) > 5000 else cleaned_text

        if len(sample_text.split()) < 10:  # Too few words for reliable detection
            return "en"

        # Use multiple detection attempts for better accuracy
        detected_langs = []
        for _ in range(3):
            try:
                detected_langs.append(langdetect.detect(sample_text))
            except:
                pass

        if not detected_langs:
            return "en"

        # Return most common detection
        from collections import Counter
        most_common = Counter(detected_langs).most_common(1)[0][0]

        # Validate detected language
        supported_langs = ['en', 'es', 'fr', 'de', 'it', 'pt', 'ru', 'zh', 'ja', 'ko', 'ar', 'hi', 'bn', 'te', 'ta', 'mr', 'gu', 'ml']
        return most_common if most_common in supported_langs else "en"

    except Exception as e:
        logging.warning(f"Language detection failed: {e}")
        return "en"

def get_language_name(lang_code: str) -> str:
    """Get full language name from language code"""
    lang_mapping = {
        'en': 'English', 'es': 'Spanish', 'fr': 'French', 'de': 'German', 'it': 'Italian',
        'pt': 'Portuguese', 'ru': 'Russian', 'zh': 'Chinese', 'ja': 'Japanese', 'ko': 'Korean',
        'ar': 'Arabic', 'hi': 'Hindi', 'bn': 'Bengali', 'te': 'Telugu', 'ta': 'Tamil',
        'mr': 'Marathi', 'gu': 'Gujarati', 'ml': 'Malayalam'
    }
    return lang_mapping.get(lang_code, 'English')

# --- URL Extraction ---
class URLExtractor:
    """Extract and validate URLs from document content"""

    @staticmethod
    def extract_urls(text: str) -> List[ExtractedURL]:
        """Extract URLs with context from text"""
        extracted_urls = []

        # Find all URLs in the text
        url_pattern = r'https?://[^\s<>"\']+|www\.[^\s<>"\']+\.[^\s<>"\']+'
        matches = re.finditer(url_pattern, text)

        for match in matches:
            url = match.group()
            if not url.startswith('http'):
                url = 'http://' + url

            # Get surrounding context
            start_idx = max(0, match.start() - 100)
            end_idx = min(len(text), match.end() + 100)
            context = text[start_idx:end_idx].strip()

            # Categorize URL
            url_type = URLExtractor._categorize_url(url, context)

            extracted_urls.append(ExtractedURL(
                url=url,
                context=context,
                source_location=f"Position {match.start()}",
                confidence=0.9,
                url_type=url_type
            ))

        return extracted_urls

    @staticmethod
    def _categorize_url(url: str, context: str) -> str:
        """Categorize URL based on URL pattern and context"""
        url_lower = url.lower()
        context_lower = context.lower()

        if any(term in url_lower for term in ['myfavouritecity', 'city']):
            return 'mission_city'
        elif any(term in url_lower for term in ['flightnumber', 'flight']):
            return 'mission_flight'
        elif any(term in url_lower for term in ['api', 'endpoint']):
            return 'api_endpoint'
        elif any(term in context_lower for term in ['click', 'link', 'visit']):
            return 'navigation'
        elif any(term in url_lower for term in ['image', 'img', 'photo', 'png', 'jpg']):
            return 'image'
        else:
            return 'general'

# --- Specialized Processors ---
class ImageOCRProcessor:
    """Specialized processor for PNG/JPEG images using Pytesseract"""

    @staticmethod
    async def process_image_file(file_content: bytes, file_path: str, request_id: str) -> List[ExtractedImage]:
        """Process images with Pytesseract OCR"""
        images = []
        if not OCR_AVAILABLE:
            logging.error("Pytesseract not available but image processing requested")
            return images

        temp_img_path = os.path.join(TEMP_FILES_PATH, f"{request_id}_{uuid.uuid4().hex}.png")

        try:
            # Load image from bytes
            pil_image = Image.open(BytesIO(file_content))

            # Convert to RGB if necessary
            if pil_image.mode != 'RGB':
                pil_image = pil_image.convert('RGB')

            # Save a temporary copy for metadata purposes
            pil_image.save(temp_img_path, 'PNG')
            width, height = pil_image.size

            # Perform OCR using pytesseract
            try:
                # Use image_to_data to get confidence scores
                ocr_data = pytesseract.image_to_data(pil_image, output_type=pytesseract.Output.DATAFRAME, lang='eng')

                # Filter out low-confidence words
                ocr_data = ocr_data[ocr_data.conf > 0]

                if not ocr_data.empty:
                    # Join words to form the text
                    final_ocr_text = " ".join(ocr_data['text'].dropna().astype(str))

                    # Calculate average confidence
                    average_confidence = ocr_data['conf'].mean() / 100.0 # Normalize to 0-1 range

                    images.append(ExtractedImage(
                        image_path=temp_img_path,
                        ocr_text=final_ocr_text.strip(),
                        metadata={
                            'source': file_path,
                            'extraction_method': 'pytesseract',
                            'image_dimensions': f"{width}x{height}",
                            'processing_timestamp': datetime.now().isoformat(),
                            'mean_confidence_score': f"{average_confidence:.2f}"
                        },
                        confidence=average_confidence
                    ))
                    logging.info(f"Successfully extracted text from image with Pytesseract: {len(final_ocr_text)} characters")
                else:
                    logging.warning("No text with sufficient confidence extracted from image using Pytesseract")

            except pytesseract.TesseractNotFoundError:
                logging.error("Tesseract is not installed or not in your PATH. OCR will not work.")
            except Exception as e:
                logging.error(f"Pytesseract OCR processing failed: {e}")

        except Exception as e:
            logging.error(f"Image processing failed for {file_path}: {e}")
        finally:
            # Cleanup temporary file
            if os.path.exists(temp_img_path):
                try:
                    os.remove(temp_img_path)
                except Exception as e:
                    logging.warning(f"Failed to cleanup temp image file: {e}")

        return images


class EnhancedXLSXTableExtractor:
    """Enhanced XLSX processor with improved table extraction and formatting"""

    @staticmethod
    def extract_tables_from_xlsx(file_content: bytes) -> List[ExtractedTable]:
        """Enhanced table extraction from XLSX with better formatting and context"""
        tables = []

        try:
            with pd.ExcelFile(BytesIO(file_content)) as excel_file:
                sheet_names = excel_file.sheet_names
                logging.info(f"Processing {len(sheet_names)} sheets from XLSX")

                # Process each sheet with enhanced logic
                for sheet_name in sheet_names:
                    try:
                        sheet_tables = EnhancedXLSXTableExtractor._process_sheet_enhanced(
                            excel_file, sheet_name
                        )
                        tables.extend(sheet_tables)

                    except Exception as e:
                        logging.warning(f"Failed to process sheet '{sheet_name}': {e}")
                        # Add error table for debugging
                        tables.append(ExtractedTable(
                            content=f"ERROR processing sheet '{sheet_name}': {str(e)}",
                            table_type='xlsx_error',
                            location=f'Sheet: {sheet_name}',
                            metadata={'error': str(e)}
                        ))

                # Add cross-sheet analysis
                if len(tables) > 1:
                    cross_sheet_analysis = EnhancedXLSXTableExtractor._analyze_cross_sheet_relationships(tables)
                    if cross_sheet_analysis:
                        tables.append(cross_sheet_analysis)

        except Exception as e:
            logging.error(f"Failed to process XLSX file: {e}")
            tables.append(ExtractedTable(
                content=f"XLSX Processing Error: {str(e)}",
                table_type='xlsx_error',
                location='File level',
                metadata={'error': str(e)}
            ))

        return tables

    @staticmethod
    def _process_sheet_enhanced(excel_file, sheet_name: str) -> List[ExtractedTable]:
        """Enhanced processing of individual sheet with multiple strategies"""
        sheet_tables = []

        try:
            # Strategy 1: Read with automatic header detection
            df_auto = pd.read_excel(excel_file, sheet_name=sheet_name)

            # Strategy 2: Read without headers for raw data analysis
            df_raw = pd.read_excel(excel_file, sheet_name=sheet_name, header=None)

            # Choose best strategy based on data quality
            df_chosen = EnhancedXLSXTableExtractor._choose_best_dataframe(df_auto, df_raw)

            if df_chosen.empty:
                return sheet_tables

            # Clean the dataframe
            df_cleaned = EnhancedXLSXTableExtractor._clean_dataframe(df_chosen)

            if df_cleaned.empty:
                return sheet_tables

            # Detect multiple tables in the sheet
            table_regions = EnhancedXLSXTableExtractor._detect_table_regions(df_cleaned)

            if not table_regions:
                # Treat entire sheet as one table
                table_regions = [(0, 0, len(df_cleaned)-1, len(df_cleaned.columns)-1)]

            # Process each detected table region
            for idx, (start_row, start_col, end_row, end_col) in enumerate(table_regions):
                try:
                    table_df = df_cleaned.iloc[start_row:end_row+1, start_col:end_col+1]

                    if table_df.empty:
                        continue

                    # Enhanced table processing
                    table_content = EnhancedXLSXTableExtractor._format_table_enhanced(
                        table_df, sheet_name, idx+1, start_row, start_col
                    )

                    # Enhanced metadata extraction
                    metadata = EnhancedXLSXTableExtractor._extract_enhanced_metadata(
                        table_df, sheet_name, idx+1
                    )

                    location = f'Sheet: {sheet_name}'
                    if len(table_regions) > 1:
                        location += f', Table: {idx+1}, Region: R{start_row+1}C{start_col+1}:R{end_row+1}C{end_col+1}'

                    sheet_tables.append(ExtractedTable(
                        content=table_content,
                        table_type='xlsx_enhanced',
                        location=location,
                        metadata=metadata
                    ))

                except Exception as e:
                    logging.warning(f"Error processing table region {idx+1} in sheet '{sheet_name}': {e}")
                    continue

        except Exception as e:
            logging.error(f"Enhanced sheet processing failed for '{sheet_name}': {e}")
            # Fallback to basic processing
            try:
                df_basic = pd.read_excel(excel_file, sheet_name=sheet_name, header=None)
                if not df_basic.empty:
                    basic_content = EnhancedXLSXTableExtractor._format_table_basic(df_basic, sheet_name)
                    sheet_tables.append(ExtractedTable(
                        content=basic_content,
                        table_type='xlsx_basic_fallback',
                        location=f'Sheet: {sheet_name} (Fallback)',
                        metadata={'processing_method': 'basic_fallback', 'original_error': str(e)}
                    ))
            except Exception as fallback_error:
                logging.error(f"Even basic processing failed for '{sheet_name}': {fallback_error}")

        return sheet_tables

    @staticmethod
    def _choose_best_dataframe(df_auto: pd.DataFrame, df_raw: pd.DataFrame) -> pd.DataFrame:
        """Choose the best dataframe based on data quality metrics"""
        if df_auto.empty and df_raw.empty:
            return df_auto
        elif df_auto.empty:
            return df_raw
        elif df_raw.empty:
            return df_auto

        # Calculate quality scores
        auto_score = EnhancedXLSXTableExtractor._calculate_dataframe_quality(df_auto)
        raw_score = EnhancedXLSXTableExtractor._calculate_dataframe_quality(df_raw)

        # Prefer auto if scores are close (header detection is usually better)
        if abs(auto_score - raw_score) < 0.1:
            return df_auto

        return df_auto if auto_score > raw_score else df_raw

    @staticmethod
    def _calculate_dataframe_quality(df: pd.DataFrame) -> float:
        """Calculate a quality score for dataframe"""
        if df.empty:
            return 0.0

        score = 0.0
        total_cells = df.shape[0] * df.shape[1]

        if total_cells == 0:
            return 0.0

        # Non-null ratio
        non_null_ratio = df.count().sum() / total_cells
        score += non_null_ratio * 0.4

        # Data type diversity (good sign)
        numeric_cols = df.select_dtypes(include=[np.number]).shape[1]
        text_cols = df.select_dtypes(include=['object']).shape[1]
        datetime_cols = df.select_dtypes(include=['datetime']).shape[1]

        type_diversity = min(1.0, (numeric_cols + text_cols + datetime_cols) / df.shape[1])
        score += type_diversity * 0.3

        # Header quality (for auto dataframe)
        if hasattr(df.columns, 'str') and df.shape[1] > 0:
            valid_headers = sum(1 for col in df.columns if isinstance(col, str) and len(str(col).strip()) > 0)
            header_quality = valid_headers / df.shape[1]
            score += header_quality * 0.3

        return min(score, 1.0)

    @staticmethod
    def _clean_dataframe(df: pd.DataFrame) -> pd.DataFrame:
        """Enhanced dataframe cleaning"""
        if df.empty:
            return df

        # Remove completely empty rows and columns
        df_cleaned = df.dropna(how='all').dropna(axis=1, how='all')

        if df_cleaned.empty:
            return df_cleaned

        # Fill NaN values with empty strings for better processing
        df_cleaned = df_cleaned.fillna('')

        # Clean string columns
        for col in df_cleaned.columns:
            if df_cleaned[col].dtype == 'object':
                df_cleaned[col] = df_cleaned[col].astype(str).str.strip()
                # Replace multiple whitespaces with single space
                df_cleaned[col] = df_cleaned[col].str.replace(r'\s+', ' ', regex=True)

        return df_cleaned

    @staticmethod
    def _detect_table_regions(df: pd.DataFrame) -> List[tuple]:
        """Detect multiple table regions within a sheet"""
        if df.empty:
            return []

        # Simple region detection: look for empty rows/columns that might separate tables
        regions = []

        # For now, implement basic single-table detection
        if not df.empty:
            regions.append((0, 0, len(df)-1, len(df.columns)-1))

        return regions

    @staticmethod
    def _format_table_enhanced(df: pd.DataFrame, sheet_name: str, table_num: int,
                             start_row: int, start_col: int) -> str:
        """Enhanced table formatting with better context and readability"""
        lines = [f"=== SHEET: {sheet_name} - TABLE {table_num} ==="]

        if df.empty:
            lines.append("EMPTY TABLE")
            return "\n".join(lines)

        # Add table position info
        lines.append(f"POSITION: Starting at Row {start_row+1}, Column {start_col+1}")
        lines.append(f"DIMENSIONS: {df.shape[0]} rows × {df.shape[1]} columns")
        lines.append("")

        # Process headers with better formatting
        headers = [str(col)[:30] for col in df.columns]
        lines.append("COLUMN HEADERS:")
        for i, header in enumerate(headers):
            lines.append(f"  Col_{i+1}: {header}")
        lines.append("")

        # Add sample data with enhanced formatting
        lines.append("DATA PREVIEW:")
        lines.append("-" * 80)

        # Show header row
        header_row = " | ".join([f"{h[:15]:15s}" for h in headers])
        lines.append(f"HEADERS: {header_row}")
        lines.append("-" * 80)

        # Show data rows (limit to prevent excessive content)
        max_rows = min(20, len(df))
        for idx in range(max_rows):
            row = df.iloc[idx]
            row_values = []
            for val in row:
                val_str = str(val)[:15] if pd.notna(val) and val != '' else ""
                row_values.append(f"{val_str:15s}")

            row_str = " | ".join(row_values)
            lines.append(f"ROW_{idx+1:2d}: {row_str}")

        if len(df) > max_rows:
            lines.append(f"... [{len(df) - max_rows} more rows]")

        # Add data analysis
        lines.append("")
        lines.append("DATA ANALYSIS:")

        # Column statistics
        for col in df.columns:
            col_data = df[col]
            non_empty = col_data[col_data != ''].count()
            unique_vals = col_data[col_data != ''].nunique()

            analysis = f"  {str(col)[:20]}: {non_empty} values, {unique_vals} unique"

            # Add sample values for better context
            sample_vals = col_data[col_data != ''].head(3).tolist()
            if sample_vals:
                sample_str = ", ".join([str(v)[:10] for v in sample_vals])
                analysis += f" (samples: {sample_str})"

            lines.append(analysis)

        # Check for mission-relevant content
        mission_indicators = EnhancedXLSXTableExtractor._detect_mission_content_enhanced(df)
        if mission_indicators:
            lines.append("")
            lines.append("MISSION CONTENT DETECTED:")
            lines.extend([f"  {indicator}" for indicator in mission_indicators])

        return "\n".join(lines)

    @staticmethod
    def _format_table_basic(df: pd.DataFrame, sheet_name: str) -> str:
        """Basic fallback table formatting"""
        lines = [f"=== SHEET: {sheet_name} (Basic Processing) ==="]

        if df.empty:
            lines.append("EMPTY SHEET")
            return "\n".join(lines)

        lines.append(f"DIMENSIONS: {df.shape[0]} rows × {df.shape[1]} columns")
        lines.append("")

        # Simple row-by-row output
        max_rows = min(15, len(df))
        for idx in range(max_rows):
            row = df.iloc[idx]
            row_str = " | ".join([str(val)[:20] if pd.notna(val) else "" for val in row])
            lines.append(f"ROW_{idx+1}: {row_str}")

        if len(df) > max_rows:
            lines.append(f"... [{len(df) - max_rows} more rows]")

        return "\n".join(lines)

    @staticmethod
    def _extract_enhanced_metadata(df: pd.DataFrame, sheet_name: str, table_num: int) -> Dict[str, Any]:
        """Extract comprehensive metadata from table"""
        if df.empty:
            return {'error': 'Empty dataframe'}

        metadata = {
            'sheet_name': sheet_name,
            'table_number': table_num,
            'dimensions': df.shape,
            'extraction_method': 'enhanced_xlsx_processing',
            'processing_timestamp': pd.Timestamp.now().isoformat()
        }

        # Data type analysis
        data_types = {}
        for col in df.columns:
            col_data = df[col][df[col] != '']
            if col_data.empty:
                data_types[str(col)] = 'empty'
            else:
                # Try to infer data type
                numeric_count = sum(1 for val in col_data if str(val).replace('.', '').replace('-', '').isdigit())
                if numeric_count / len(col_data) > 0.8:
                    data_types[str(col)] = 'numeric'
                elif any(keyword in str(col).lower() for keyword in ['date', 'time', 'created', 'modified']):
                    data_types[str(col)] = 'datetime'
                else:
                    data_types[str(col)] = 'text'

        metadata['column_types'] = data_types

        # Content analysis
        total_cells = df.shape[0] * df.shape[1]
        non_empty_cells = sum(1 for col in df.columns for val in df[col] if val != '')

        metadata['data_density'] = non_empty_cells / total_cells if total_cells > 0 else 0
        metadata['non_empty_cells'] = non_empty_cells

        # Mission content detection
        metadata['contains_mission_data'] = EnhancedXLSXTableExtractor._check_for_mission_data_enhanced(df)

        # URL detection
        all_text = ' '.join([str(val) for col in df.columns for val in df[col] if val != ''])
        url_pattern = r'https?://[^\s<>"\']+|www\.[^\s<>"\']+\.[^\s<>"\']+'
        urls_found = re.findall(url_pattern, all_text)
        metadata['urls_found'] = len(urls_found)
        metadata['contains_urls'] = len(urls_found) > 0

        return metadata

    @staticmethod
    def _detect_mission_content_enhanced(df: pd.DataFrame) -> List[str]:
        """Enhanced mission content detection"""
        indicators = []

        if df.empty:
            return indicators

        # Check column names
        column_text = ' '.join([str(col).lower() for col in df.columns])
        mission_keywords = ['city', 'flight', 'landmark', 'token', 'api', 'url', 'mission', 'favourite']

        for keyword in mission_keywords:
            if keyword in column_text:
                indicators.append(f"Mission keyword '{keyword}' found in column headers")

        # Check data content
        all_data = []
        for col in df.columns:
            for val in df[col]:
                if pd.notna(val) and val != '':
                    all_data.append(str(val).lower())

        content_text = ' '.join(all_data)

        # Look for specific mission patterns
        patterns = {
            'city': r'\b[A-Za-z]+(?:\s+[A-Za-z]+)*(?:\s+city)?\b',
            'flight': r'\b[A-Z]{1,3}\d{3,4}\b',
            'url': r'https?://[^\s<>"\']+',
            'api': r'\bapi\b|\bendpoint\b',
            'token': r'\btoken\b|\bsecret\b'
        }

        for pattern_name, pattern in patterns.items():
            matches = re.findall(pattern, content_text)
            if matches:
                indicators.append(f"Found {len(matches)} {pattern_name} pattern(s): {matches[:3]}...")

        return indicators

    @staticmethod
    def _check_for_mission_data_enhanced(df: pd.DataFrame) -> bool:
        """Enhanced check for mission-related data"""
        if df.empty:
            return False

        indicators = EnhancedXLSXTableExtractor._detect_mission_content_enhanced(df)
        return len(indicators) > 0

    @staticmethod
    def _analyze_cross_sheet_relationships(tables: List[ExtractedTable]) -> ExtractedTable:
        """Analyze relationships between different sheets"""
        try:
            analysis_parts = ["=== CROSS-SHEET ANALYSIS ==="]

            # Extract sheet names
            sheet_names = []
            for table in tables:
                if 'sheet_name' in table.metadata:
                    sheet_names.append(table.metadata['sheet_name'])

            if len(set(sheet_names)) > 1:
                analysis_parts.append(f"WORKBOOK CONTAINS: {len(set(sheet_names))} sheets")
                analysis_parts.append(f"SHEET NAMES: {', '.join(set(sheet_names))}")
                analysis_parts.append("")

                # Look for common themes across sheets
                all_content = ' '.join([table.content for table in tables])

                # Mission content analysis across sheets
                if 'mission' in all_content.lower() or 'flight' in all_content.lower():
                    analysis_parts.append("CROSS-SHEET MISSION CONTENT DETECTED")
                    analysis_parts.append("This workbook may contain related mission data across multiple sheets")
                    analysis_parts.append("")

                # URL analysis across sheets
                url_pattern = r'https?://[^\s<>"\']+|www\.[^\s<>"\']+\.[^\s<>"\']+'
                all_urls = re.findall(url_pattern, all_content)
                if all_urls:
                    analysis_parts.append(f"URLS FOUND ACROSS SHEETS: {len(all_urls)} total")
                    unique_urls = list(set(all_urls))
                    analysis_parts.extend([f"  {url}" for url in unique_urls[:5]])
                    if len(unique_urls) > 5:
                        analysis_parts.append(f"  ... and {len(unique_urls) - 5} more")

                return ExtractedTable(
                    content='\n'.join(analysis_parts),
                    table_type='xlsx_cross_analysis',
                    location='Cross-sheet analysis',
                    metadata={
                        'analysis_type': 'cross_sheet_relationships',
                        'sheets_analyzed': list(set(sheet_names)),
                        'total_urls_found': len(all_urls) if all_urls else 0
                    }
                )

        except Exception as e:
            logging.warning(f"Cross-sheet analysis failed: {e}")

        return None


class EnhancedPPTXTextExtractor:
    """Enhanced PPTX processor with improved text and table extraction"""

    @staticmethod
    def extract_text_from_pptx(file_content: bytes) -> str:
        """Enhanced text extraction from PPTX with better structure and table handling"""
        text_parts = []

        try:
            presentation = Presentation(BytesIO(file_content))

            # Extract presentation metadata
            if hasattr(presentation.core_properties, 'title') and presentation.core_properties.title:
                text_parts.append(f"=== PRESENTATION TITLE ===")
                text_parts.append(presentation.core_properties.title)
                text_parts.append("")

            total_slides = len(presentation.slides)
            logging.info(f"Processing {total_slides} slides from PPTX")

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = EnhancedPPTXTextExtractor._extract_slide_content(slide, slide_num)
                if slide_content:
                    text_parts.extend(slide_content)
                    text_parts.append("")  # Add spacing between slides

            # Add presentation summary
            text_parts.append(f"=== PRESENTATION SUMMARY ===")
            text_parts.append(f"Total slides processed: {total_slides}")

            logging.info(f"Successfully extracted enhanced content from {total_slides} PPTX slides")

        except Exception as e:
            logging.error(f"Enhanced PPTX extraction failed: {e}")
            text_parts.append(f"Error processing PPTX: {str(e)}")

        return "\n".join(text_parts)

    @staticmethod
    def _extract_slide_content(slide, slide_num: int) -> List[str]:
        """Extract comprehensive content from a single slide"""
        slide_parts = [f"=== SLIDE {slide_num} ==="]

        # Track different content types
        titles = []
        bullet_points = []
        paragraphs = []
        table_contents = []
        urls = []

        # Process all shapes in the slide
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # Handle text-containing shapes
                if hasattr(shape, "text") and shape.text.strip():
                    content = EnhancedPPTXTextExtractor._process_text_shape(shape)
                    if content:
                        shape_type = EnhancedPPTXTextExtractor._determine_shape_type(shape)

                        if shape_type == "title":
                            titles.append(content)
                        elif shape_type == "bullets":
                            bullet_points.extend(content.split('\n'))
                        else:
                            paragraphs.append(content)

                # Handle table shapes
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_content = EnhancedPPTXTextExtractor._extract_table_from_slide(
                        shape.table, f"Slide_{slide_num}_Table_{shape_idx+1}"
                    )
                    if table_content:
                        table_contents.append(table_content)

                # Handle group shapes (nested content)
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_content = EnhancedPPTXTextExtractor._extract_group_content(shape)
                    if group_content:
                        paragraphs.append(group_content)

            except Exception as e:
                logging.warning(f"Error processing shape {shape_idx} in slide {slide_num}: {e}")
                continue

        # Organize content by type
        if titles:
            slide_parts.append("SLIDE TITLE:")
            slide_parts.extend([f"  {title}" for title in titles])
            slide_parts.append("")

        if bullet_points:
            slide_parts.append("KEY POINTS:")
            for point in bullet_points:
                point = point.strip()
                if point:
                    # Preserve bullet structure
                    if point.startswith(('•', '▪', '◦', '-', '*')):
                        slide_parts.append(f"  {point}")
                    else:
                        slide_parts.append(f"  • {point}")
            slide_parts.append("")

        if paragraphs:
            slide_parts.append("CONTENT:")
            for para in paragraphs:
                slide_parts.append(f"  {para}")
            slide_parts.append("")

        if table_contents:
            slide_parts.append("TABLES:")
            slide_parts.extend(table_contents)
            slide_parts.append("")

        # Extract and process slide notes
        notes_content = EnhancedPPTXTextExtractor._extract_slide_notes(slide)
        if notes_content:
            slide_parts.append("SLIDE NOTES:")
            slide_parts.append(f"  {notes_content}")
            slide_parts.append("")

        # Extract URLs from all slide content
        all_slide_text = " ".join(slide_parts)
        urls = EnhancedPPTXTextExtractor._extract_urls_from_text(all_slide_text)
        if urls:
            slide_parts.append("URLS FOUND:")
            for url in urls:
                slide_parts.append(f"  {url}")
            slide_parts.append("")

        return slide_parts if len(slide_parts) > 1 else []

    @staticmethod
    def _process_text_shape(shape) -> str:
        """Process text from shape with better formatting"""
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return shape.text.strip()

        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            para_text = paragraph.text.strip()
            if para_text:
                # Preserve indentation and bullet structure
                level = paragraph.level if hasattr(paragraph, 'level') else 0
                indent = "  " * level

                # Check for bullet points
                if any(para_text.startswith(marker) for marker in ['•', '▪', '◦', '-', '*']):
                    text_parts.append(f"{indent}{para_text}")
                else:
                    # Add bullet if it's in a list context and doesn't have one
                    if level > 0:
                        text_parts.append(f"{indent}• {para_text}")
                    else:
                        text_parts.append(para_text)

        return "\n".join(text_parts)

    @staticmethod
    def _determine_shape_type(shape) -> str:
        """Determine the semantic type of a shape"""
        if not hasattr(shape, 'text'):
            return "unknown"

        text = shape.text.strip()

        # Check for title characteristics
        if (len(text.split()) <= 10 and
            (text.isupper() or
             (hasattr(shape, 'top') and hasattr(shape, 'left') and
              shape.top < 1000000))):  # Top area of slide
            return "title"

        # Check for bullet points
        if any(marker in text for marker in ['•', '▪', '◦', '‣']) or text.count('\n') > 2:
            return "bullets"

        return "content"

    @staticmethod
    def _extract_table_from_slide(table, table_id: str) -> str:
        """Extract and format table content from slide"""
        try:
            table_parts = [f"--- {table_id} ---"]

            rows = []
            headers = []

            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_data.append(cell_text)

                if row_idx == 0:
                    headers = row_data
                    table_parts.append(f"HEADERS: {' | '.join(headers)}")
                    table_parts.append("-" * min(80, len(' | '.join(headers))))
                else:
                    rows.append(row_data)
                    row_str = ' | '.join([str(cell)[:25] for cell in row_data])
                    table_parts.append(f"ROW_{row_idx}: {row_str}")

            table_parts.append(f"TABLE SUMMARY: {len(rows)} rows, {len(headers)} columns")
            return "\n".join(table_parts)

        except Exception as e:
            logging.warning(f"Error extracting table {table_id}: {e}")
            return f"Error extracting table {table_id}: {str(e)}"

    @staticmethod
    def _extract_group_content(group_shape) -> str:
        """Extract content from grouped shapes"""
        group_texts = []
        try:
            for shape in group_shape.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    group_texts.append(shape.text.strip())
        except Exception as e:
            logging.warning(f"Error processing group shape: {e}")

        return " | ".join(group_texts) if group_texts else ""

    @staticmethod
    def _extract_slide_notes(slide) -> str:
        """Extract and clean slide notes"""
        try:
            if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        # Clean up notes text
                        notes_text = re.sub(r'\s+', ' ', notes_text)
                        return notes_text
        except Exception as e:
            logging.warning(f"Error extracting slide notes: {e}")

        return ""

    @staticmethod
    def _extract_urls_from_text(text: str) -> List[str]:
        """Extract unique URLs from text"""
        url_pattern = r'https?://[^\s<>"\']+|www\.[^\s<>"\']+\.[^\s<>"\']+'
        urls = re.findall(url_pattern, text)
        # Remove duplicates while preserving order
        seen = set()
        unique_urls = []
        for url in urls:
            if url not in seen:
                seen.add(url)
                unique_urls.append(url)
        return unique_urls


# --- PDF Processing ---
class EnhancedPDFProcessor:
    """Enhanced PDF processor with better text extraction"""

    @staticmethod
    def extract_pdf_content(doc_content: bytes) -> str:
        """Extract text from PDF with layout preservation"""
        text_parts = []

        try:
            with fitz.open(stream=BytesIO(doc_content), filetype="pdf") as pdf_doc:
                # Extract document metadata
                metadata = pdf_doc.metadata
                if metadata:
                    meta_parts = []
                    if metadata.get("title"):
                        meta_parts.append(f"Title: {metadata['title']}")
                    if metadata.get("author"):
                        meta_parts.append(f"Author: {metadata['author']}")
                    if metadata.get("subject"):
                        meta_parts.append(f"Subject: {metadata['subject']}")

                    if meta_parts:
                        text_parts.append("=== DOCUMENT METADATA ===")
                        text_parts.append("\n".join(meta_parts))
                        text_parts.append("")

                # Process each page
                for page_num, page in enumerate(pdf_doc, 1):
                    page_texts = [f"=== PAGE {page_num} ==="]

                    # Get page text with better formatting
                    text_dict = page.get_text("dict")
                    blocks = text_dict.get("blocks", [])

                    page_text_blocks = []
                    for block in blocks:
                        if block.get("type") == 0:  # Text block
                            # Process text blocks
                            for line in block.get("lines", []):
                                line_text = " ".join(span.get("text", "") for span in line.get("spans", []))
                                if line_text.strip():
                                    page_text_blocks.append(line_text)

                    if page_text_blocks:
                        page_texts.append("\n".join(page_text_blocks))
                    else:
                        # Fallback to simpler extraction
                        page_text = page.get_text("text")
                        if page_text.strip():
                            page_texts.append(page_text)

                    # Extract images (for debugging, we don't process them here)
                    image_list = page.get_images(full=True)
                    if image_list:
                        page_texts.append(f"[Page contains {len(image_list)} images]")

                    # Extract links
                    links = page.get_links()
                    urls = []
                    for link in links:
                        if link.get("uri") and link["uri"].startswith("http"):
                            urls.append(link["uri"])

                    if urls:
                        page_texts.append(f"URLS: {', '.join(urls)}")

                    if len(page_texts) > 1:  # More than just the page header
                        text_parts.extend(page_texts)
                        text_parts.append("")  # Add spacing between pages

                logging.info(f"Extracted text from {len(pdf_doc)} PDF pages")

        except Exception as e:
            logging.error(f"PDF extraction failed: {e}")
            text_parts.append(f"Error extracting PDF content: {str(e)}")

        return "\n".join(text_parts)

# --- Enhanced Document Processor ---
class TargetedDocumentProcessor:
    """Targeted document processor with strict type handling"""

    @staticmethod
    async def process_document(doc_content: bytes, doc_type: str, doc_url: str, request_id: str) -> ProcessedDocument:
        """Process documents with type-specific handlers"""
        text = ""
        tables = []
        images = []
        extracted_urls = []

        loop = asyncio.get_running_loop()

        try:
            with ThreadPoolExecutor(max_workers=2) as executor:
                if doc_type == "pdf":
                    text = await loop.run_in_executor(
                        executor,
                        EnhancedPDFProcessor.extract_pdf_content,
                        doc_content
                    )

                elif doc_type == "docx":
                    text = await loop.run_in_executor(
                        executor,
                        TargetedDocumentProcessor._extract_docx_text,
                        doc_content
                    )

                elif doc_type == "pptx":
                    text = await loop.run_in_executor(
                        executor,
                        EnhancedPPTXTextExtractor.extract_text_from_pptx,
                        doc_content
                    )

                elif doc_type == "xlsx":
                    tables = await loop.run_in_executor(
                        executor,
                        EnhancedXLSXTableExtractor.extract_tables_from_xlsx,
                        doc_content
                    )
                    text = "Spreadsheet document with structured data tables."

                elif doc_type in ["png", "jpeg"]:
                    images = await ImageOCRProcessor.process_image_file(
                        doc_content,
                        doc_url,
                        request_id
                    )
                    text = "Image document processed with OCR."

                elif doc_type == "html":
                    try:
                        docs = await loop.run_in_executor(
                            executor,
                            lambda: WebBaseLoader(web_paths=[doc_url]).load()
                        )
                        if docs:
                            text = docs[0].page_content
                        else:
                            # Fallback direct HTML processing
                            response = await loop.run_in_executor(
                                executor,
                                lambda: requests.get(doc_url, timeout=30)
                            )
                            soup = BeautifulSoup(response.text, 'html.parser')
                            for script in soup(["script", "style", "nav", "footer", "header"]):
                                script.decompose()
                            text = soup.get_text(separator=' ', strip=True)
                    except Exception as e:
                        logging.error(f"Failed to process HTML: {e}")
                        text = f"Error processing HTML: {str(e)}"

                elif doc_type == "txt":
                    try:
                        text = doc_content.decode('utf-8', errors='replace')
                    except UnicodeDecodeError:
                        # Try different encodings
                        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                            try:
                                text = doc_content.decode(encoding, errors='replace')
                                break
                            except:
                                pass

                else:
                    text = f"Unsupported document type: {doc_type}"
                    logging.warning(f"Unsupported document type: {doc_type}")

        except Exception as e:
            logging.error(f"Document processing failed for type {doc_type}: {e}")
            text = f"Error processing {doc_type} document: {str(e)}"

        # Extract URLs from text
        extracted_urls = URLExtractor.extract_urls(text)

        # Combine all content
        enhanced_content = text

        # Add table content
        if tables:
            table_content = "\n\n".join([
                f"--- TABLE FROM {table.location} ---\n{table.content}"
                for table in tables
            ])
            enhanced_content += f"\n\n=== EXTRACTED TABLES ===\n{table_content}"

        # Add image content
        if images:
            image_content = "\n\n".join([
                f"--- IMAGE: {img.metadata.get('source', 'unknown')} ---\n{img.ocr_text}"
                for img in images
            ])
            enhanced_content += f"\n\n=== TEXT FROM IMAGES (OCR) ===\n{image_content}"

        # Add URL summary
        if extracted_urls:
            url_content = "\n".join([
                f"URL: {url.url} (Type: {url.url_type})"
                for url in extracted_urls
            ])
            enhanced_content += f"\n\n=== EXTRACTED URLS ===\n{url_content}"

        # Detect language from the complete content
        detected_language = detect_language_robust(enhanced_content)
        language_name = get_language_name(detected_language)
        logging.info(f"Detected document language: {language_name} ({detected_language})")

        return ProcessedDocument(
            content=enhanced_content,
            metadata={
                "doc_type": doc_type,
                "doc_url": doc_url,
                "char_count": len(enhanced_content),
                "table_count": len(tables),
                "image_count": len(images),
                "url_count": len(extracted_urls),
                "language": language_name,
                "language_code": detected_language,
                "processing_timestamp": datetime.now().isoformat()
            },
            tables=tables,
            images=images,
            extracted_urls=extracted_urls,
            detected_language=detected_language
        )

    @staticmethod
    def _extract_docx_text(doc_content: bytes) -> str:
        """Enhanced DOCX text extraction"""
        text_parts = []

        try:
            doc = docx.Document(BytesIO(doc_content))

            # Extract document properties if available
            core_properties = doc.core_properties
            if core_properties:
                prop_parts = []
                if core_properties.title:
                    prop_parts.append(f"Title: {core_properties.title}")
                if core_properties.author:
                    prop_parts.append(f"Author: {core_properties.author}")
                if core_properties.subject:
                    prop_parts.append(f"Subject: {core_properties.subject}")

                if prop_parts:
                    text_parts.append("=== DOCUMENT PROPERTIES ===")
                    text_parts.extend(prop_parts)
                    text_parts.append("")

            # Extract heading-based structure
            current_heading = None
            paragraph_buffer = []

            for para in doc.paragraphs:
                if not para.text.strip():
                    continue

                # Check if this is a heading
                if para.style.name.startswith('Heading'):
                    # Process previous section
                    if paragraph_buffer:
                        if current_heading:
                            text_parts.append(f"=== {current_heading} ===")
                        text_parts.append("\n".join(paragraph_buffer))
                        text_parts.append("")
                        paragraph_buffer = []

                    current_heading = para.text.strip()
                else:
                    paragraph_buffer.append(para.text.strip())

            # Process final section
            if paragraph_buffer:
                if current_heading:
                    text_parts.append(f"=== {current_heading} ===")
                text_parts.append("\n".join(paragraph_buffer))

            # Extract table content
            for table_idx, table in enumerate(doc.tables, 1):
                table_text = []
                table_text.append(f"=== TABLE {table_idx} ===")

                for row_idx, row in enumerate(table.rows):
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        prefix = "HEADER: " if row_idx == 0 else f"ROW {row_idx}: "
                        table_text.append(f"{prefix}{row_text}")

                if len(table_text) > 1:  # More than just the table header
                    text_parts.append("\n".join(table_text))
                    text_parts.append("")

            logging.info(f"Extracted text from DOCX: {len(text_parts)} elements")

        except Exception as e:
            logging.error(f"DOCX extraction failed: {e}")
            text_parts.append(f"Error extracting DOCX content: {str(e)}")

        return "\n".join(text_parts)

# --- Enhanced Chunking Strategy ---
class AdaptiveChunkingStrategy:
    """Adaptive chunking based on document characteristics"""

    @staticmethod
    def create_adaptive_chunks(processed_doc: ProcessedDocument) -> List[Document]:
        """Create adaptive chunks based on document characteristics"""
        content = processed_doc.content
        doc_type = processed_doc.metadata.get("doc_type", "unknown")
        has_tables = len(processed_doc.tables) > 0
        has_images = len(processed_doc.images) > 0

        # Determine optimal chunk parameters
        chunk_params = AdaptiveChunkingStrategy._determine_chunk_params(
            doc_type=doc_type,
            has_tables=has_tables,
            has_images=has_images,
            content_length=len(content)
        )

        # Create splitter with adaptive parameters
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_params["chunk_size"],
            chunk_overlap=chunk_params["chunk_overlap"],
            separators=chunk_params["separators"],
            length_function=len,
            is_separator_regex=False,
            keep_separator=True
        )

        # Split content
        chunks = splitter.create_documents([content])

        # Enhance chunks with metadata
        enhanced_chunks = []
        for i, chunk in enumerate(chunks):
            # Calculate importance score
            importance_score = AdaptiveChunkingStrategy._calculate_importance_score(
                chunk.page_content,
                doc_type=doc_type,
                has_tables="=== EXTRACTED TABLES ===" in chunk.page_content,
                has_images="=== TEXT FROM IMAGES ===" in chunk.page_content
            )

            # Detect content type
            content_type = AdaptiveChunkingStrategy._detect_content_type(chunk.page_content)

            # Create enhanced metadata
            chunk.metadata = {
                "chunk_id": i,
                "char_count": len(chunk.page_content),
                "word_count": len(chunk.page_content.split()),
                "has_tables": "=== EXTRACTED TABLES ===" in chunk.page_content,
                "has_images": "=== TEXT FROM IMAGES ===" in chunk.page_content,
                "has_urls": "=== EXTRACTED URLS ===" in chunk.page_content or bool(re.search(r'https?://\S+', chunk.page_content)),
                "importance_score": importance_score,
                "content_type": content_type,
                "doc_type": doc_type
            }

            enhanced_chunks.append(chunk)

        logging.info(f"Created {len(enhanced_chunks)} adaptive chunks (params: {chunk_params})")
        return enhanced_chunks

    @staticmethod
    def _determine_chunk_params(doc_type: str, has_tables: bool, has_images: bool, content_length: int) -> Dict:
        """Determine optimal chunking parameters based on document characteristics"""

        # Base parameters
        params = {
            "chunk_size": 2000,
            "chunk_overlap": 400,
            "separators": ["\n\n\n", "\n\n", "\n", ". ", "? ", "! ", "; ", ", "]
        }

        # Adjust for document type
        if doc_type == "xlsx" or has_tables:
            # Larger chunks for structured data
            params["chunk_size"] = 3000
            params["chunk_overlap"] = 600
            params["separators"] = ["\n=== EXTRACTED TABLES ===", "\n--- TABLE", "\n\n\n", "\n\n", "\n", ". "]

        elif doc_type in ["pdf", "docx"] and content_length > 50000:
            # Medium-large chunks for long documents
            params["chunk_size"] = 2500
            params["chunk_overlap"] = 500

        elif doc_type in ["png", "jpeg"] or has_images:
            # Smaller chunks for OCR content
            params["chunk_size"] = 1500
            params["chunk_overlap"] = 300
            params["separators"] = ["\n=== TEXT FROM IMAGES ===", "\n--- IMAGE", "\n\n", "\n", ". "]

        elif doc_type == "html":
            # Medium chunks for HTML
            params["chunk_size"] = 2200
            params["chunk_overlap"] = 400

        elif doc_type == "pptx":
            # Medium chunks for presentations with slide breaks
            params["chunk_size"] = 2500
            params["chunk_overlap"] = 500
            params["separators"] = ["\n=== SLIDE", "\n\n\n", "\n\n", "\n", ". "]

        # Adjust for content length
        if content_length < 10000:
            # Smaller chunks for very short documents
            params["chunk_size"] = min(params["chunk_size"], content_length // 2)
            params["chunk_overlap"] = params["chunk_size"] // 5

        elif content_length > 100000:
            # Larger chunks for very long documents
            params["chunk_size"] = 3000
            params["chunk_overlap"] = 600

        return params

    @staticmethod
    def _calculate_importance_score(content: str, doc_type: str, has_tables: bool, has_images: bool) -> float:
        """Calculate importance score based on content characteristics"""
        score = 0.5  # Base score

        # Boost for document structure markers
        if any(marker in content for marker in ["===", "---", "TABLE", "SHEET", "SLIDE"]):
            score += 0.15

        # Boost for mission-critical content
        mission_terms = ["mission", "flight", "city", "landmark", "token", "api", "url", "favourite"]
        mission_term_count = sum(1 for term in mission_terms if term in content.lower())
        if mission_term_count > 0:
            score += min(0.25, 0.05 * mission_term_count)

        # Boost for URLs
        if re.search(r'https?://\S+', content):
            score += 0.2

        # Boost for numerical data
        if re.search(r'\b\d+\.\d+\b|\b\d{2,}\b', content):
            score += 0.1

        # Boost for tables
        if has_tables:
            score += 0.15

        # Boost for OCR text
        if has_images:
            score += 0.15

        return min(score, 1.0)

    @staticmethod
    def _detect_content_type(content: str) -> str:
        """Detect the type of content in the chunk"""
        content_lower = content.lower()

        if "=== EXTRACTED TABLES ===" in content:
            return "table_data"
        elif "=== TEXT FROM IMAGES ===" in content:
            return "ocr_text"
        elif "=== EXTRACTED URLS ===" in content or re.search(r'https?://\S+', content):
            return "url_content"
        elif any(term in content_lower for term in ["mission", "brief", "objective", "city", "landmark", "flight"]):
            return "mission_content"
        elif re.search(r'(\d+\.\d+|\d{4,})', content):
            return "numerical_data"
        else:
            return "general_text"

# --- Enhanced Retrieval System ---
class EnhancedRetriever:
    """Enhanced retriever with re-ranking capabilities"""

    def __init__(self, chunks: List[Document], embedding_model, use_reranking: bool = True):
        self.chunks = chunks
        self.embedding_model = embedding_model
        self.use_reranking = use_reranking and RERANK_AVAILABLE and reranker is not None
        self.vectorstore = None
        self._build_vectorstore()

    def _build_vectorstore(self):
        """Build vector store with enhanced parameters"""
        if not self.chunks:
            return

        try:
            # Create FAISS vectorstore
            self.vectorstore = FAISS.from_documents(
                documents=self.chunks,
                embedding=self.embedding_model
            )
            logging.info(f"Built vectorstore with {len(self.chunks)} chunks")
        except Exception as e:
            logging.error(f"Failed to build vectorstore: {e}")
            self.vectorstore = None

    async def get_relevant_documents(self, query: str, k: int = 15) -> List[Document]:
        """Enhanced document retrieval with re-ranking"""
        if not self.vectorstore:
            return []

        try:
            # Initial retrieval with higher k for re-ranking
            initial_k = min(k * 3, len(self.chunks)) if self.use_reranking else k

            # Use similarity search with scores
            docs_with_scores = await self.vectorstore.asimilarity_search_with_relevance_scores(
                query,
                k=initial_k
            )

            # Extract documents and scores
            candidate_docs = [doc for doc, _ in docs_with_scores]

            if not candidate_docs:
                logging.warning("No documents retrieved from vector search")
                return []

            # Re-rank if available
            if self.use_reranking and len(candidate_docs) > 3:
                reranked_docs = await self._rerank_documents(query, candidate_docs, k)
                logging.info(f"Re-ranked {len(candidate_docs)} -> {len(reranked_docs)} documents")
                return reranked_docs
            else:
                # Fallback to importance-based filtering
                return self._importance_filter(candidate_docs, k)

        except Exception as e:
            logging.error(f"Enhanced retrieval failed: {e}")
            # Fallback to basic similarity search
            try:
                return await self.vectorstore.asimilarity_search(query, k=min(k, len(self.chunks)))
            except Exception as e:
                logging.error(f"Fallback retrieval failed: {e}")
                return []

    async def _rerank_documents(self, query: str, documents: List[Document], k: int) -> List[Document]:
        """Re-rank documents using cross-encoder"""
        try:
            # Prepare query-document pairs
            pairs = [(query, doc.page_content) for doc in documents]

            # Get re-ranking scores
            scores = reranker.predict(pairs)

            # Combine documents with scores
            scored_docs = list(zip(documents, scores))

            # Sort by re-ranking score (descending)
            scored_docs.sort(key=lambda x: x[1], reverse=True)

            # Get top k documents
            top_docs = [doc for doc, _ in scored_docs[:k]]

            return top_docs

        except Exception as e:
            logging.error(f"Re-ranking failed: {e}")
            return self._importance_filter(documents, k)

    def _importance_filter(self, documents: List[Document], k: int) -> List[Document]:
        """Filter documents by importance score and content diversity"""
        # Score documents by importance
        scored_docs = [(doc, doc.metadata.get("importance_score", 0.5)) for doc in documents]

        # Sort by importance score (descending)
        scored_docs.sort(key=lambda x: x[1], reverse=True)

        # Simple diversity filter - avoid consecutive chunks from same section
        diverse_docs = []
        prev_content_type = None

        # First get highest importance docs regardless of diversity
        top_docs = [doc for doc, _ in scored_docs[:k//2]]
        diverse_docs.extend(top_docs)

        # Then add remaining docs with diversity consideration
        for doc, _ in scored_docs[k//2:]:
            content_type = doc.metadata.get("content_type", "general_text")

            if len(diverse_docs) >= k:
                break

            # Add document if content type is different from previous
            if content_type != prev_content_type:
                diverse_docs.append(doc)
                prev_content_type = content_type

        # Fill remaining slots with highest importance docs not yet included
        if len(diverse_docs) < k:
            remaining_docs = [doc for doc, _ in scored_docs if doc not in diverse_docs]
            diverse_docs.extend(remaining_docs[:k - len(diverse_docs)])

        return diverse_docs[:k]

# --- Mission Detection and Execution ---
def detect_mission_document(doc: ProcessedDocument) -> bool:
    """Enhanced mission document detection"""
    if not doc or not doc.content:
        return False

    content_lower = doc.content.lower()

    # Primary mission indicators
    primary_keywords = [
        "sachin's parallel world",
        "mission brief",
        "myfavouritecity",
        "getfirstcityflightnumber"
    ]

    # Secondary mission indicators
    secondary_keywords = [
        "flight number",
        "landmark",
        "city",
        "api endpoint",
        "secret token"
    ]

    # URL patterns
    url_patterns = [
        "myfavouritecity",
        "getfirstcityflightnumber",
        "flightnumber"
    ]

    primary_matches = sum(1 for keyword in primary_keywords if keyword in content_lower)
    secondary_matches = sum(1 for keyword in secondary_keywords if keyword in content_lower)
    url_matches = sum(1 for pattern in url_patterns if pattern in content_lower)

    # Mission document if:
    # - 2+ primary keywords, OR
    # - 1 primary + 3+ secondary keywords, OR
    # - 2+ URL patterns
    is_mission = (
        primary_matches >= 2 or
        (primary_matches >= 1 and secondary_matches >= 3) or
        url_matches >= 2
    )

    logging.info(f"Mission detection - Primary: {primary_matches}, Secondary: {secondary_matches}, URLs: {url_matches}, Is Mission: {is_mission}")
    return is_mission

def create_mission_solving_agent_prompt(detected_language: str) -> ChatPromptTemplate:
    """Create mission-solving agent prompt with language enforcement"""
    language_name = get_language_name(detected_language).upper()

    system_message = f"""**CRITICAL: ALL RESPONSES MUST BE IN {language_name} ({detected_language}). NON-NEGOTIABLE.**

You are a mission execution specialist. Your task is to find the final flight number by following these exact steps:

**MISSION EXECUTION PROTOCOL:**
1. **STEP 1**: Locate the URL containing `myFavouriteCity` in the provided context
2. **STEP 2**: Use `fetch_contextual_url_content` tool to call that URL and get the city name
3. **STEP 3**: Find the city name in the tables/context to identify its corresponding landmark
4. **STEP 4**: Locate the URL with the landmark pattern (get...FlightNumber) in the context
5. **STEP 5**: Use `fetch_contextual_url_content` tool to call the final URL and get flight number

**OUTPUT REQUIREMENT:**
- Your final response must contain ONLY the flight number
- No additional text, explanations, or formatting
- Example: "AI101" or "6E2045" (just the flight number)
- The flight number MUST be in {language_name}"""

    # The 'agent_scratchpad' placeholder is essential for the agent to remember previous steps
    return ChatPromptTemplate.from_messages([
        ("system", system_message),
        ("human", "Execute the mission protocol. CONTEXT:\n---\n{context}\n---"),
        MessagesPlaceholder(variable_name="agent_scratchpad"),
    ])

# --- Language-Strict Response Generation ---
def generate_language_strict_prompt(detected_language: str) -> str:
    """Generate prompt with strict language enforcement"""
    language_name = get_language_name(detected_language).upper()

    # A generic "not found" message in different languages
    not_found_messages = {
        'en': "This information is not available in the provided document context.",
        'es': "Esta información no está disponible en el contexto del documento proporcionado.",
        'fr': "Cette information n'est pas disponible dans le contexte du document fourni.",
        'de': "Diese Information ist im bereitgestellten Dokumentkontext nicht verfügbar.",
        'it': "Questa informazione non è disponibile nel contesto del documento fornito.",
        'hi': "यह जानकारी प्रदान किए गए दस्तावेज़ के संदर्भ में उपलब्ध नहीं है।",
        'bn': "প্রদত্ত নথির প্রেক্ষাপটে এই তথ্য উপলব্ধ নয়।",
        'te': "సమర్పించిన పత్రం సందర్భంలో ఈ సమాచారం అందుబాటులో లేదు.",
        'ta': "வழங்கப்பட்ட ஆவண சூழலில் இந்த தகவல் கிடைக்கவில்லை.",
        'mr': "दिलेल्या दस्तऐवजाच्या संदर्भात ही माहिती उपलब्ध नाही.",
    }

    not_found_message = not_found_messages.get(detected_language, not_found_messages['en'])

    return f"""**ABSOLUTE, UNBREAKABLE RULE: YOUR FINAL ANSWER MUST BE IN {language_name} ({detected_language}). THIS IS YOUR PRIMARY OBJECTIVE.**

You are an expert document analyst. Your task is to answer the user's question, which may be in a different language from the context provided.

**ANALYSIS PROTOCOL:**
1.  **Understand the Question:** First, fully understand the user's question here: `{{input}}`.
2.  **Find the Answer:** Next, locate the answer to that question within the following context, which is in **{language_name}**.
3.  **Formulate the Response:** Finally, formulate your answer *exclusively* in **{language_name}**, using only information found in the context.

---
**CONTEXT (in {language_name}):**
{{context}}
---

**USER QUESTION (Language may vary):** `{{input}}`

**CRITICAL DIRECTIVES:**
- **DO NOT refuse to answer because the languages are different.** Your specific purpose is to bridge this language gap.
- If the answer exists in the context, provide it in **{language_name}**.
- If the answer truly does not exist in the context, you MUST respond with the exact phrase: "{not_found_message}"
- Keep answers concise (under 80 words) but informative.
- Never mention that you're translating or that languages are different.

**Answer in {language_name} (max 80 words):**"""

# --- Enhanced Agent Tools ---
@tool
async def fetch_contextual_url_content(url: str, context_hint: str = "") -> str:
    """Enhanced URL fetching with better error handling"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'application/json,text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.5',
            'Accept-Encoding': 'gzip, deflate',
            'Connection': 'keep-alive',
            'Cache-Control': 'no-cache',
        }

        response = requests.get(url, timeout=45, headers=headers)
        response.raise_for_status()

        content_type = response.headers.get('Content-Type', '').lower()

        if 'application/json' in content_type:
            try:
                json_data = response.json()
                return json.dumps(json_data, indent=2, ensure_ascii=False)
            except json.JSONDecodeError:
                return response.text
        elif 'text/html' in content_type:
            soup = BeautifulSoup(response.text, 'html.parser')

            # Remove unwanted elements
            for element in soup(["script", "style", "nav", "footer", "header", "aside"]):
                element.decompose()

            # Get clean text
            text = soup.get_text(separator=' ', strip=True)
            text = re.sub(r'\s+', ' ', text).strip()

            return text
        else:
            return response.text

    except requests.RequestException as e:
        error_msg = f"Failed to fetch URL {url}: {str(e)}"
        logging.error(error_msg)
        return error_msg

# --- Enhanced Core Processing Logic ---
async def process_single_question(question: str, index: int, enhanced_retriever: EnhancedRetriever,
                                  mission_executor: AgentExecutor, doc: ProcessedDocument,
                                  is_mission: bool, logger: logging.Logger) -> Tuple[int, str]:
    """Process a single question with enhanced error handling"""
    try:
        if is_mission:
            logger.info(f"Q{index + 1}: Executing mission protocol")
            if index == 0:
                # First question executes the mission
                response = await mission_executor.ainvoke({"context": doc.content})
                mission_result = response.get("output", "Mission execution failed").strip()

                # Extract just the flight number if response contains extra text
                flight_pattern = r'\b[A-Z0-9]{2,3}[\s-]?\d{3,4}\b'
                flight_match = re.search(flight_pattern, mission_result)
                if flight_match:
                    return index, flight_match.group().replace(" ", "").replace("-", "")

                return index, mission_result
            # For subsequent questions in mission document
            return index, "Mission objective completed. Refer to first answer."

        # Standard RAG processing
        logger.info(f"Q{index + 1}: Executing enhanced RAG protocol")
        if not enhanced_retriever:
            logger.error("No retriever available")
            not_found_messages = {
                'en': "This information is not available in the provided document context.",
                'es': "Esta información no está disponible en el contexto del documento proporcionado.",
                'fr': "Cette information n'est pas disponible dans le contexte du document fourni.",
                'hi': "यह जानकारी प्रदान किए गए दस्तावेज़ के संदर्भ में उपलब्ध नहीं है।",
                'ta': "வழங்கப்பட்ட ஆவண சூழலில் இந்த தகவல் கிடைக்கவில்லை.",
            }
            return index, not_found_messages.get(doc.detected_language, not_found_messages['en'])

        # Enhanced document retrieval with adaptive k
        query_complexity = len(question.split())
        adaptive_k = min(12 if query_complexity > 10 else 8, len(enhanced_retriever.chunks))

        docs = await enhanced_retriever.get_relevant_documents(question, k=adaptive_k)

        if not docs:
            logger.warning(f"Q{index + 1}: No relevant documents found")
            not_found_messages = {
                'en': "This information is not available in the provided document context.",
                'es': "Esta información no está disponible en el contexto del documento proporcionado.",
                'fr': "Cette information n'est pas disponible dans le contexte du document fourni.",
                'hi': "यह जानकारी प्रदान किए गए दस्तावेज़ के संदर्भ में उपलब्ध नहीं है।",
                'ta': "வழங்கப்பட்ட ஆவண சூழலில் இந்த தகவல் கிடைக்கவில்லை.",
            }
            return index, not_found_messages.get(doc.detected_language, not_found_messages['en'])

        # Log retrieval stats
        logger.info(f"Q{index + 1}: Retrieved {len(docs)} documents with importance scores: {[d.metadata.get('importance_score', 0.5) for d in docs[:3]]}")

        # Language-strict prompt with better context formatting
        prompt = ChatPromptTemplate.from_template(generate_language_strict_prompt(doc.detected_language))
        chain = create_stuff_documents_chain(llm, prompt)

        # Format context with metadata for better processing
        formatted_docs = []
        for i, doc_chunk in enumerate(docs):
            importance = doc_chunk.metadata.get('importance_score', 0.5)
            content_type = doc_chunk.metadata.get('content_type', 'general')
            formatted_docs.append(f"[Source {i+1} - Relevance: {importance:.2f} - Type: {content_type}]\n{doc_chunk.page_content}")

        enhanced_context_docs = [Document(page_content="\n\n".join(formatted_docs))]

        # Invoke chain with question and context
        result = await chain.ainvoke({"context": enhanced_context_docs, "input": question})

        # Clean and validate result for language consistency
        cleaned_result = clean_and_validate_answer(result, doc.detected_language)

        logger.info(f"Q{index + 1}: Generated answer ({len(cleaned_result.split())} words)")
        return index, cleaned_result

    except Exception as e:
        logger.error(f"Processing error for Q{index+1}: {e}", exc_info=True)
        return index, f"An error occurred during processing. ({str(e)})"

def clean_and_validate_answer(result: str, detected_language: str) -> str:
    """Clean and validate answer, ensuring language consistency"""
    if not result or len(result.strip()) < 3:
        # Return not found message in the detected language
        not_found_messages = {
            'en': "This information is not available in the provided document context.",
            'es': "Esta información no está disponible en el contexto del documento proporcionado.",
            'fr': "Cette information n'est pas disponible dans le contexte du document fourni.",
            'de': "Diese Information ist im bereitgestellten Dokumentkontext nicht verfügbar.",
            'it': "Questa informazione non è disponibile nel contesto del documento fornito.",
            'hi': "यह जानकारी प्रदान किए गए दस्तावेज़ के संदर्भ में उपलब्ध नहीं है।",
            'bn': "প্রদত্ত নথির প্রেক্ষাপটে এই তথ্য উপলব্ধ নয়।",
            'te': "సమర్పించిన పత్రం సందర్భంలో ఈ సమాచారం అందుబాటులో లేదు.",
            'ta': "வழங்கப்பட்ட ஆவண சூழலில் இந்த தகவல் கிடைக்கவில்லை.",
            'mr': "दिलेल्या दस्तऐवजाच्या संदर्भात ही माहिती उपलब्ध नाही.",
        }
        return not_found_messages.get(detected_language, not_found_messages['en'])

    answer = result.strip()

    # Remove common AI-generated prefixes while preserving non-English content
    if detected_language == "en":
        prefixes_to_remove = [
            r'^Based on the (?:enhanced )?context(?:,?)\s*',
            r'^According to the (?:provided )?(?:enhanced )?(?:context|document)(?:,?)\s*',
            r'^The (?:enhanced )?context (?:shows|indicates) that\s*',
            r'^From the (?:provided )?(?:enhanced )?(?:context|sources)(?:,?)\s*'
        ]
        for prefix_pattern in prefixes_to_remove:
            answer = re.sub(prefix_pattern, '', answer, flags=re.IGNORECASE)

    # Clean up whitespace
    answer = re.sub(r'\s+', ' ', answer).strip()

    # Enforce word limit (80 words max)
    words = answer.split()
    if len(words) > 80:
        # Try to keep complete sentences within limit
        sentences = re.split(r'[.!?।？！]', answer)
        truncated_sentences = []
        word_count = 0

        for sentence in sentences:
            sentence_words = sentence.strip().split()
            if word_count + len(sentence_words) <= 80:
                truncated_sentences.append(sentence.strip())
                word_count += len(sentence_words)
            else:
                break

        if truncated_sentences:
            answer = '. '.join(truncated_sentences)
            if not answer.endswith(('.', '!', '?', '।', '？', '！')):
                answer += '.'
        else:
            # Fallback: hard truncate at 80 words
            answer = ' '.join(words[:80]) + '...'

    # Add proper ending punctuation based on language
    if answer and not answer.endswith(('.', '!', '?', '।', '？', '！', '...')):
        if detected_language in ['hi', 'bn', 'te', 'ta', 'mr', 'gu']:
            answer += '।'
        else:
            answer += '.'

    # Capitalize first letter for Latin-script languages
    if answer and detected_language in ['en', 'es', 'fr', 'de', 'it', 'pt'] and answer[0].islower():
        answer = answer[0].upper() + answer[1:]

    return answer

# --- Main Processing Function ---
async def process_query_enhanced(doc_url: str, questions: List[str], logger: logging.Logger, request_id: str) -> List[str]:
    """Main processing function with enhanced error handling"""
    logger.info(f"Starting processing for {len(questions)} questions")

    try:
        # Download and process document
        response = requests.get(doc_url, timeout=120)
        response.raise_for_status()

        # Detect document type
        doc_type = detect_document_type_strict(doc_url)
        logger.info(f"Processing document type: {doc_type}")

        # Process document with appropriate handler
        processed_doc = await TargetedDocumentProcessor.process_document(
            response.content,
            doc_type,
            doc_url,
            request_id
        )

        logger.info(f"Document processed: {len(processed_doc.content)} chars, {len(processed_doc.tables)} tables, "
                   f"{len(processed_doc.images)} images, {len(processed_doc.extracted_urls)} URLs")
        logger.info(f"Detected language: {processed_doc.metadata.get('language')} ({processed_doc.detected_language})")

        # Check if this is a mission document
        is_mission_doc = detect_mission_document(processed_doc)
        logger.info(f"Mission document detected: {is_mission_doc}")

    except Exception as e:
        logger.error(f"Document processing failed: {e}", exc_info=True)
        return [f"Document processing failed: {str(e)}"] * len(questions)

    # Set up mission executor if needed
    mission_executor = None
    if is_mission_doc:
        try:
            tools = [fetch_contextual_url_content]
            mission_prompt = create_mission_solving_agent_prompt(processed_doc.detected_language)
            mission_agent = create_tool_calling_agent(llm, tools, mission_prompt)
            mission_executor = AgentExecutor(
                agent=mission_agent,
                tools=tools,
                verbose=True,
                handle_parsing_errors=True,
                max_iterations=5  # Increased for complex missions
            )
            logger.info("Mission executor set up successfully")
        except Exception as e:
            logger.error(f"Failed to set up mission executor: {e}", exc_info=True)
            is_mission_doc = False

    # Set up enhanced retriever for RAG
    enhanced_retriever = None
    if not is_mission_doc and processed_doc.content:
        try:
            logger.info("Creating adaptive chunks...")
            # Create adaptive chunks
            chunks = AdaptiveChunkingStrategy.create_adaptive_chunks(processed_doc)

            if chunks:
                # Select appropriate embedding model
                doc_size = len(processed_doc.content)
                has_complex_content = len(processed_doc.tables) > 0 or len(processed_doc.images) > 0

                # Use more accurate embedding model for complex content or large documents
                embedding_model = embeddings_accurate if (has_complex_content or doc_size > 30000) else embeddings_fast
                if not embedding_model:
                    embedding_model = embeddings_fast or embeddings_accurate

                if not embedding_model:
                    raise ValueError("No embedding model available")

                logger.info(f"Using {'accurate' if embedding_model == embeddings_accurate else 'fast'} embeddings")

                # Create enhanced retriever
                enhanced_retriever = EnhancedRetriever(
                    chunks=chunks,
                    embedding_model=embedding_model,
                    use_reranking=RERANK_AVAILABLE
                )

                logger.info(f"Enhanced retriever created with {len(chunks)} chunks")

        except Exception as e:
            logger.error(f"Failed to set up retriever: {e}", exc_info=True)
            enhanced_retriever = None

    # Process questions concurrently
    tasks = []
    for i, question in enumerate(questions):
        task = process_single_question(
            question,
            i,
            enhanced_retriever,
            mission_executor,
            processed_doc,
            is_mission_doc,
            logger
        )
        tasks.append(task)

    results = await asyncio.gather(*tasks, return_exceptions=True)

    # Process results with enhanced error handling
    ordered_answers = ["Processing failed."] * len(questions)
    successful = 0

    for result in results:
        if isinstance(result, Exception):
            logger.error(f"Task failed: {result}", exc_info=True)
        elif isinstance(result, tuple) and len(result) == 2:
            idx, answer = result
            ordered_answers[idx] = answer
            successful += 1
        else:
            logger.error(f"Invalid result format: {result}")

    logger.info(f"Processing completed: {successful}/{len(questions)} questions processed successfully")

    return ordered_answers

# --- API Endpoints ---
@app.post("/api/v1/hackrx/run", response_model=QueryResponse, dependencies=[Depends(validate_token)])
async def run_enhanced_submission(request: QueryRequest):
    """Enhanced API endpoint for document processing"""
    request_id = str(uuid.uuid4())
    logger = setup_request_logger(request_id)

    try:
        logger.info("="*70)
        logger.info(f"START PROCESSING - Request ID: {request_id}")
        logger.info(f"Document URL: {request.documents}")
        logger.info(f"Questions: {len(request.questions)} received")

        start_time = asyncio.get_event_loop().time()

        # Process with enhanced system
        answers = await process_query_enhanced(
            request.documents,
            request.questions,
            logger,
            request_id
        )

        end_time = asyncio.get_event_loop().time()
        processing_time = end_time - start_time

        logger.info(f"Processing completed in {processing_time:.2f} seconds")
        logger.info(f"Average time per question: {processing_time/len(request.questions):.2f}s")

        return QueryResponse(answers=answers)

    except HTTPException as e:
        logger.error(f"HTTP Error: {e.detail}")
        raise e
    except Exception as e:
        logger.error(f"Unexpected error: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"System error: {str(e)}")
    finally:
        # Cleanup logger
        handlers = logger.handlers[:]
        for handler in handlers:
            handler.close()
            logger.removeHandler(handler)
        logging.info(f"END Request ID: {request_id}")

@app.get("/health")
async def health_check():
    """Enhanced health check with feature status"""
    return {
        "status": "healthy",
        "version": "13.0 - Enhanced Language-Strict Document-Targeted RAG with Improved PPTX/XLSX Processing",
        "features": {
            "adaptive_chunking": True,
            "targeted_processing": True,
            "enhanced_pptx_processing": True,
            "enhanced_xlsx_processing": True,
            "reranking": RERANK_AVAILABLE,
            "ocr_multilingual": OCR_AVAILABLE,
            "language_detection": LANG_DETECT_AVAILABLE,
            "mission_detection": True
        },
        "models": {
            "embedding_fast": "sentence-transformers/all-MiniLM-L6-v2" if embeddings_fast else "Not loaded",
            "embedding_accurate": "BAAI/bge-small-en-v1.5" if embeddings_accurate else "Not loaded",
            "reranker": "cross-encoder/ms-marco-MiniLM-L-6-v2" if reranker else "Not available",
            "llm": "gemini-2.5-flash-lite" if llm else "Not loaded",
            "ocr": "pytesseract" if OCR_AVAILABLE else "Not available"
        },
        "supported_formats": {
            "documents": ["pdf", "docx", "pptx", "txt", "html"],
            "tables": ["xlsx"],
            "images": ["png", "jpeg", "jpg"]
        },
        "language_support": {
            "output_languages": ["English", "Spanish", "French", "German", "Italian", "Portuguese",
                               "Hindi", "Bengali", "Telugu", "Tamil", "Marathi", "Malayalam"]
        },
        "device": device,
        "timestamp": datetime.now().isoformat()
    }

@app.get("/api/v1/capabilities")
async def list_capabilities():
    """List detailed system capabilities"""
    return {
        "document_processing": {
            "pdf": {
                "handler": "EnhancedPDFProcessor",
                "features": ["text extraction", "metadata extraction", "URL detection", "page structure preservation"]
            },
            "docx": {
                "handler": "TargetedDocumentProcessor._extract_docx_text",
                "features": ["text extraction", "table extraction", "structure preservation", "heading detection"]
            },
            "pptx": {
                "handler": "EnhancedPPTXTextExtractor",
                "features": ["slide content extraction", "table extraction", "notes extraction", "URL detection", "bullet point processing", "structured content organization"]
            },
            "xlsx": {
                "handler": "EnhancedXLSXTableExtractor",
                "features": ["advanced table extraction", "multi-sheet analysis", "header detection", "mission data detection", "cross-sheet relationships", "enhanced formatting"]
            },
            "images": {
                "handler": "ImageOCRProcessor",
                "features": ["OCR with Pytesseract", "confidence scores", "mission data detection"]
            }
        },
        "retrieval": {
            "handler": "EnhancedRetriever",
            "features": ["vector similarity", "cross-encoder reranking", "importance filtering", "diversity promotion"]
        },
        "chunking": {
            "handler": "AdaptiveChunkingStrategy",
            "features": ["document-type-specific parameters", "content-aware chunking", "metadata enrichment"]
        },
        "language": {
            "detection": LANG_DETECT_AVAILABLE,
            "strict_enforcement": True,
            "question_language": "Any",
            "answer_language": "Same as document"
        },
        "mission": {
            "detection": True,
            "execution": True,
            "agent_tools": ["fetch_contextual_url_content"]
        },
        "improvements": {
            "pptx": [
                "Enhanced slide content extraction with better structure",
                "Table extraction from slides",
                "Improved bullet point and hierarchical content handling",
                "Better slide notes processing",
                "Comprehensive URL extraction",
                "Group shape content processing"
            ],
            "xlsx": [
                "Multiple processing strategies with quality-based selection",
                "Enhanced table formatting with detailed analysis",
                "Cross-sheet relationship analysis",
                "Mission content detection with pattern matching",
                "Data type inference and statistics",
                "Comprehensive metadata extraction"
            ]
        }
    }

if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=int(os.getenv("PORT", 8000)), workers=1)
