
import os
import re
import uuid
import logging
import shutil
import asyncio
from io import BytesIO
from typing import List, Dict, Tuple, Any
from datetime import datetime
from PIL import Image

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from data_models import ExtractedImage
from config import TEMP_FILES_PATH, OCR_AVAILABLE

# Import pytesseract conditionally to handle environments where it's not available
if OCR_AVAILABLE:
    import pytesseract
    
# --- ENHANCED PPTX PROCESSOR ---
class AdvancedPPTXProcessor:
    """Advanced PPTX processor with comprehensive slide content extraction"""

    @staticmethod
    async def process_pptx(file_content: bytes, request_id: str) -> Tuple[str, List[ExtractedImage]]:
        """Process PPTX with enhanced extraction of text, images, and relationships"""
        text_parts = []
        extracted_images = []
        slide_references = {}  # Track references between slides

        temp_dir = os.path.join(TEMP_FILES_PATH, request_id)
        os.makedirs(temp_dir, exist_ok=True)

        try:
            presentation = Presentation(BytesIO(file_content))

            # Extract presentation metadata
            metadata = AdvancedPPTXProcessor._extract_presentation_metadata(presentation)
            if metadata:
                text_parts.extend(metadata)

            total_slides = len(presentation.slides)
            logging.info(f"Processing {total_slides} slides from PPTX")

            # Process slides in two passes
            # First pass: Extract all content and record references
            slide_contents = {}
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content, references, images = await AdvancedPPTXProcessor._process_slide_comprehensive(
                    slide, slide_num, temp_dir, request_id
                )
                slide_contents[slide_num] = slide_content
                slide_references[slide_num] = references
                extracted_images.extend(images)

            # Second pass: Resolve cross-references
            for slide_num, content in slide_contents.items():
                resolved_content = AdvancedPPTXProcessor._resolve_cross_references(
                    content, slide_references, slide_contents
                )
                text_parts.extend(resolved_content)
                text_parts.append("")  # Add spacing between slides

            # Add presentation summary with relationship map
            summary = AdvancedPPTXProcessor._generate_presentation_summary(
                total_slides, slide_references, slide_contents
            )
            text_parts.extend(summary)

            logging.info(f"Successfully extracted advanced content from {total_slides} PPTX slides")

        except Exception as e:
            logging.error(f"Advanced PPTX extraction failed: {e}", exc_info=True)
            text_parts.append(f"Error processing PPTX: {str(e)}")

        finally:
            # Clean up temp files
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except Exception as e:
                logging.warning(f"Failed to clean up temp directory: {e}")

        return "\n".join(text_parts), extracted_images

    @staticmethod
    def _extract_presentation_metadata(presentation) -> List[str]:
        """Extract comprehensive metadata from presentation"""
        metadata_parts = []

        try:
            core_props = presentation.core_properties
            metadata_parts.append("=== PRESENTATION METADATA ===")

            if hasattr(core_props, 'title') and core_props.title:
                metadata_parts.append(f"Title: {core_props.title}")

            if hasattr(core_props, 'subject') and core_props.subject:
                metadata_parts.append(f"Subject: {core_props.subject}")

            if hasattr(core_props, 'author') and core_props.author:
                metadata_parts.append(f"Author: {core_props.author}")

            if hasattr(core_props, 'comments') and core_props.comments:
                metadata_parts.append(f"Comments: {core_props.comments}")

            if hasattr(core_props, 'keywords') and core_props.keywords:
                metadata_parts.append(f"Keywords: {core_props.keywords}")

            if hasattr(core_props, 'last_modified_by') and core_props.last_modified_by:
                metadata_parts.append(f"Last modified by: {core_props.last_modified_by}")

            if hasattr(core_props, 'revision') and core_props.revision:
                metadata_parts.append(f"Revision: {core_props.revision}")

            # Add presentation-wide statistics
            metadata_parts.append(f"Total slides: {len(presentation.slides)}")
            metadata_parts.append("")

        except Exception as e:
            logging.warning(f"Error extracting presentation metadata: {e}")

        return metadata_parts

    @staticmethod
    async def _process_slide_comprehensive(slide, slide_num: int, temp_dir: str, request_id: str) -> Tuple[List[str], List[str], List[ExtractedImage]]:
        """Process a slide comprehensively extracting text, images, and references"""
        slide_parts = [f"=== SLIDE {slide_num} ==="]
        references = []
        extracted_images = []

        # Detect slide type/layout
        slide_layout_name = slide.slide_layout.name if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name') else "Unknown"
        slide_parts.append(f"SLIDE TYPE: {slide_layout_name}")

        # Process different content types
        try:
            # Extract titles
            titles = AdvancedPPTXProcessor._extract_slide_titles(slide)
            if titles:
                slide_parts.append("SLIDE TITLE:")
                slide_parts.extend([f"  {title}" for title in titles])
                slide_parts.append("")

            # Extract text content and bullet points
            text_content = AdvancedPPTXProcessor._extract_slide_text_content(slide)
            if text_content:
                slide_parts.extend(text_content)
                slide_parts.append("")

            # Extract tables
            tables = AdvancedPPTXProcessor._extract_slide_tables(slide, slide_num)
            if tables:
                slide_parts.append("TABLES:")
                slide_parts.extend(tables)
                slide_parts.append("")

            # Process images
            image_results = await AdvancedPPTXProcessor._extract_slide_images(slide, slide_num, temp_dir, request_id)
            image_descriptions, slide_images = image_results

            if image_descriptions:
                slide_parts.append("IMAGES:")
                slide_parts.extend(image_descriptions)
                slide_parts.append("")

            extracted_images.extend(slide_images)

            # Process charts
            chart_descriptions = AdvancedPPTXProcessor._extract_slide_charts(slide, slide_num)
            if chart_descriptions:
                slide_parts.append("CHARTS:")
                slide_parts.extend(chart_descriptions)
                slide_parts.append("")

            # Process SmartArt
            smartart_text = AdvancedPPTXProcessor._extract_slide_smartart(slide)
            if smartart_text:
                slide_parts.append("SMARTART:")
                slide_parts.extend(smartart_text)
                slide_parts.append("")

            # Extract speaker notes
            notes = AdvancedPPTXProcessor._extract_slide_notes(slide)
            if notes:
                slide_parts.append("SPEAKER NOTES:")
                slide_parts.extend([f"  {note}" for note in notes])
                slide_parts.append("")

            # Extract hyperlinks and cross-references
            hyperlinks, slide_refs = AdvancedPPTXProcessor._extract_slide_links(slide, slide_num)

            if hyperlinks:
                slide_parts.append("HYPERLINKS:")
                slide_parts.extend(hyperlinks)
                slide_parts.append("")

            references.extend(slide_refs)

            # Look for explicit references to other slides in the content
            content_text = " ".join(slide_parts)
            ref_matches = re.finditer(r'(?:see|refer to|go to|reference)\s+(?:slide|page)?\s*(?:number)?\s*(\d+)', content_text, re.IGNORECASE)
            for match in ref_matches:
                referenced_slide = int(match.group(1))
                if 1 <= referenced_slide <= 1000:  # Reasonable slide number limit
                    references.append(f"slide_{referenced_slide}")

        except Exception as e:
            slide_parts.append(f"ERROR processing slide content: {str(e)}")
            logging.error(f"Error processing slide {slide_num}: {e}")

        return slide_parts, references, extracted_images

    @staticmethod
    def _extract_slide_titles(slide) -> List[str]:
        """Extract all title elements from a slide"""
        titles = []

        for shape in slide.shapes:
            if hasattr(shape, "is_title") and shape.is_title:
                if hasattr(shape, "text") and shape.text.strip():
                    titles.append(shape.text.strip())
            elif hasattr(shape, "text") and shape.text.strip():
                # Check if shape contains title-like text (uppercase, limited words)
                text = shape.text.strip()
                if len(text) <= 100 and (text.isupper() or text.istitle()):
                    if not titles:  # Only add as potential title if we don't have one already
                        titles.append(text)

        return titles

    @staticmethod
    def _extract_slide_text_content(slide) -> List[str]:
        """Extract formatted text content including bullet points"""
        content_parts = []
        bullet_points = []
        paragraphs = []

        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue

            if hasattr(shape, "is_title") and shape.is_title:
                continue  # Skip titles as they're handled separately

            if not hasattr(shape, "text_frame"):
                paragraphs.append(shape.text.strip())
                continue

            # Process text with paragraph formatting
            for para in shape.text_frame.paragraphs:
                if not para.text.strip():
                    continue

                level = para.level if hasattr(para, "level") else 0
                indent = "  " * level

                # Check if this is likely a bullet point
                if hasattr(para, "bullet") and para.bullet:
                    bullet_points.append(f"{indent}• {para.text.strip()}")
                elif level > 0:
                    # Treat indented paragraphs as bullet points
                    bullet_points.append(f"{indent}• {para.text.strip()}")
                else:
                    paragraphs.append(para.text.strip())

        # Add bullet points if any
        if bullet_points:
            content_parts.append("KEY POINTS:")
            content_parts.extend(bullet_points)

        # Add paragraphs if any
        if paragraphs:
            content_parts.append("CONTENT:")
            content_parts.extend([f"  {p}" for p in paragraphs])

        return content_parts

    @staticmethod
    def _extract_slide_tables(slide, slide_num: int) -> List[str]:
        """Extract tables from a slide with enhanced formatting"""
        tables = []
        table_count = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_count += 1

                try:
                    table_parts = [f"--- Table {slide_num}.{table_count} ---"]

                    rows = []
                    headers = []
                    max_col_widths = []

                    # First pass: calculate column widths
                    for row_idx, row in enumerate(shape.table.rows):
                        if row_idx == 0:
                            max_col_widths = [len(cell.text.strip()) + 2 for cell in row.cells]
                        else:
                            for col_idx, cell in enumerate(row.cells):
                                if col_idx < len(max_col_widths):
                                    max_col_widths[col_idx] = max(max_col_widths[col_idx], len(cell.text.strip()) + 2)

                    # Set reasonable maximum width
                    max_col_widths = [min(w, 30) for w in max_col_widths]

                    # Second pass: format rows
                    for row_idx, row in enumerate(shape.table.rows):
                        row_data = []
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if col_idx < len(max_col_widths):
                                width = max_col_widths[col_idx]
                                row_data.append(f"{cell_text:{width}}")
                            else:
                                row_data.append(cell_text)

                        if row_idx == 0:
                            headers = row_data
                            header_str = "| " + " | ".join(headers) + " |"
                            table_parts.append(header_str)
                            table_parts.append("|-" + "-|-".join(["-" * len(h) for h in headers]) + "-|")
                        else:
                            rows.append(row_data)
                            table_parts.append("| " + " | ".join(row_data) + " |")

                    tables.extend(table_parts)

                except Exception as e:
                    tables.append(f"Error extracting table {table_count}: {str(e)}")

        return tables

    @staticmethod
    async def _extract_slide_images(slide, slide_num: int, temp_dir: str, request_id: str) -> Tuple[List[str], List[ExtractedImage]]:
        """Extract and process images from a slide"""
        image_descriptions = []
        extracted_images = []
        image_count = 0

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1

                try:
                    # Save image to temp file
                    img_filename = f"slide_{slide_num}_img_{image_count}_{request_id}.png"
                    img_path = os.path.join(temp_dir, img_filename)

                    # Extract image
                    if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                        with open(img_path, "wb") as f:
                            f.write(shape.image.blob)

                        # Process image with OCR if available
                        if OCR_AVAILABLE:
                            with Image.open(img_path) as pil_image:
                                width, height = pil_image.size

                                # Perform OCR
                                ocr_text = pytesseract.image_to_string(pil_image)
                                ocr_text = ocr_text.strip()

                                if ocr_text:
                                    image_descriptions.append(f"  Image {image_count}: {width}x{height} - Contains text: {ocr_text[:100]}")

                                    # Create extracted image object
                                    extracted_image = ExtractedImage(
                                        image_path=img_path,
                                        ocr_text=ocr_text,
                                        metadata={
                                            'source': f"Slide {slide_num}",
                                            'extraction_method': 'pytesseract',
                                            'image_dimensions': f"{width}x{height}",
                                            'position': f"Shape {shape_idx}",
                                            'processing_timestamp': datetime.now().isoformat()
                                        },
                                        confidence=0.7  # Default confidence
                                    )
                                    extracted_images.append(extracted_image)
                                else:
                                    image_descriptions.append(f"  Image {image_count}: {width}x{height} - No text detected")
                        else:
                            image_descriptions.append(f"  Image {image_count}: Image extraction only (OCR not available)")

                except Exception as e:
                    image_descriptions.append(f"  Error extracting image {image_count}: {str(e)}")
                    logging.error(f"Error extracting image from slide {slide_num}: {e}")

        if image_count == 0:
            return [], []

        return image_descriptions, extracted_images

    @staticmethod
    def _extract_slide_charts(slide, slide_num: int) -> List[str]:
        """Extract and describe charts from a slide"""
        chart_descriptions = []
        chart_count = 0

        for shape in slide.shapes:
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart_count += 1
                chart = shape.chart

                try:
                    chart_type = str(chart.chart_type).split('.')[-1] if hasattr(chart, "chart_type") else "Unknown"

                    description = f"  Chart {chart_count}: Type: {chart_type}"

                    # Try to extract chart title
                    if hasattr(chart, "chart_title") and hasattr(chart.chart_title, "text_frame"):
                        title = chart.chart_title.text_frame.text.strip()
                        if title:
                            description += f", Title: {title}"

                    # Try to extract categories and series names
                    categories = []
                    series_names = []

                    if hasattr(chart, "plots") and chart.plots:
                        plot = chart.plots[0]
                        if hasattr(plot, "categories") and plot.categories:
                            categories = [cat.label.text if hasattr(cat, "label") and hasattr(cat.label, "text") else str(cat)
                                          for cat in plot.categories]

                        if hasattr(plot, "series"):
                            series_names = [series.name if hasattr(series, "name") else f"Series {i}"
                                            for i, series in enumerate(plot.series)]

                    if categories:
                        description += f"\n    Categories: {', '.join(categories[:5])}"
                        if len(categories) > 5:
                            description += f" and {len(categories) - 5} more"

                    if series_names:
                        description += f"\n    Series: {', '.join(series_names[:5])}"
                        if len(series_names) > 5:
                            description += f" and {len(series_names) - 5} more"

                    chart_descriptions.append(description)

                except Exception as e:
                    chart_descriptions.append(f"  Error analyzing chart {chart_count}: {str(e)}")

        return chart_descriptions

    @staticmethod
    def _extract_slide_smartart(slide) -> List[str]:
        """Extract text from SmartArt graphics"""
        smartart_text = []

        for shape in slide.shapes:
            if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # SmartArt is often grouped shapes
                try:
                    group_text = []
                    if hasattr(shape, "shapes"):
                        for subshape in shape.shapes:
                            if hasattr(subshape, "text") and subshape.text.strip():
                                group_text.append(subshape.text.strip())

                    if group_text:
                        smartart_text.append("  SmartArt content:")
                        for i, text in enumerate(group_text, 1):
                            smartart_text.append(f"    Item {i}: {text}")
                except Exception as e:
                    smartart_text.append(f"  Error extracting SmartArt: {str(e)}")

        return smartart_text

    @staticmethod
    def _extract_slide_notes(slide) -> List[str]:
        """Extract comprehensive speaker notes"""
        notes = []

        try:
            if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
                if slide.notes_slide and hasattr(slide.notes_slide, "notes_text_frame"):
                    notes_frame = slide.notes_slide.notes_text_frame

                    if hasattr(notes_frame, "paragraphs"):
                        for para in notes_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                notes.append(text)
                    else:
                        notes_text = notes_frame.text.strip()
                        if notes_text:
                            # Split by newlines while preserving meaningful structure
                            for line in notes_text.split('\n'):
                                clean_line = line.strip()
                                if clean_line:
                                    notes.append(clean_line)
        except Exception as e:
            notes.append(f"Error extracting notes: {str(e)}")

        return notes

    @staticmethod
    def _extract_slide_links(slide, slide_num: int) -> Tuple[List[str], List[str]]:
        """Extract hyperlinks and cross-references from a slide"""
        hyperlinks = []
        slide_references = []

        # Process shapes with hyperlinks
        for shape in slide.shapes:
            if hasattr(shape, "click_action") and shape.click_action:
                action = shape.click_action

                # Check for hyperlink
                if hasattr(action, "hyperlink") and action.hyperlink.address:
                    link_text = shape.text.strip() if hasattr(shape, "text") else "Link"
                    hyperlinks.append(f"  {link_text} → {action.hyperlink.address}")

                # Check for slide links
                if hasattr(action, "target_slide") and action.target_slide:
                    target_slide_idx = action.target_slide.slide_id
                    target_slide_num = None

                    # Try to find the slide number by ID
                    if hasattr(slide.part, "package") and hasattr(slide.part.package, "presentation"):
                        for i, s in enumerate(slide.part.package.presentation.slides, 1):
                            if hasattr(s, "slide_id") and s.slide_id == target_slide_idx:
                                target_slide_num = i
                                break

                    link_text = shape.text.strip() if hasattr(shape, "text") else "Slide link"
                    if target_slide_num:
                        hyperlinks.append(f"  {link_text} → Slide {target_slide_num}")
                        slide_references.append(f"slide_{target_slide_num}")
                    else:
                        hyperlinks.append(f"  {link_text} → Internal slide link")

        # Look for hyperlinks in text
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if hasattr(run, "hyperlink") and run.hyperlink.address:
                        hyperlinks.append(f"  {run.text.strip()} → {run.hyperlink.address}")

        return hyperlinks, slide_references

    @staticmethod
    def _resolve_cross_references(slide_content: List[str], references: Dict[int, List[str]], all_contents: Dict[int, List[str]]) -> List[str]:
        """Enhance slide content by resolving cross-references"""
        # For now, just add a references section if there are any
        if not references:
            return slide_content

        slide_num = int(slide_content[0].split()[-1]) if slide_content and "SLIDE" in slide_content[0] else 0
        if slide_num not in references or not references[slide_num]:
            return slide_content

        ref_slides = []
        for ref in references[slide_num]:
            if ref.startswith("slide_"):
                ref_num = int(ref.split("_")[1])
                ref_slides.append(ref_num)

        if ref_slides:
            slide_content.append("REFERENCES OTHER SLIDES:")
            for ref_slide in sorted(ref_slides):
                # Extract the title of the referenced slide
                ref_title = "Unknown Title"
                if ref_slide in all_contents:
                    for line in all_contents[ref_slide]:
                        if line.startswith("SLIDE TITLE:") and len(all_contents[ref_slide]) > all_contents[ref_slide].index(line) + 1:
                            ref_title = all_contents[ref_slide][all_contents[ref_slide].index(line) + 1].strip()
                            break

                slide_content.append(f"  Slide {ref_slide}: {ref_title}")

            slide_content.append("")

        return slide_content

        @staticmethod
        def _generate_presentation_summary(total_slides: int, slide_references: Dict[int, List[str]], slide_contents: Dict[int, List[str]]) -> List[str]:
            """Generate a presentation summary with slide relationships"""
            summary_parts = ["=== PRESENTATION SUMMARY ==="]
            summary_parts.append(f"Total slides: {total_slides}")

            # Create a map of which slides reference others
            slide_map = {}
            for slide_num, refs in slide_references.items():
                referenced_slides = []
                for ref in refs:
                    if ref.startswith("slide_"):
                        ref_num = int(ref.split("_")[1])
                        if 1 <= ref_num <= total_slides:
                            referenced_slides.append(ref_num)

                if referenced_slides:
                    slide_map[slide_num] = referenced_slides

            # Find key slides (most referenced or referencing)
            if slide_map:
                # Count how many times each slide is referenced
                reference_counts = {}
                for slide_num in range(1, total_slides + 1):
                    reference_counts[slide_num] = 0

                for slide_num, refs in slide_map.items():
                    for ref_num in refs:
                        reference_counts[ref_num] = reference_counts.get(ref_num, 0) + 1

                # Find most referenced slides
                most_referenced = sorted(reference_counts.items(), key=lambda x: x[1], reverse=True)[:3]
                if any(count > 0 for _, count in most_referenced):
                    summary_parts.append("\nKEY SLIDES:")
                    for slide_num, count in most_referenced:
                        if count > 0:
                            # Get slide title
                            title = "Unknown Title"
                            if slide_num in slide_contents:
                                for line in slide_contents[slide_num]:
                                    if line.startswith("SLIDE TITLE:") and len(slide_contents[slide_num]) > slide_contents[slide_num].index(line) + 1:
                                        title = slide_contents[slide_num][slide_contents[slide_num].index(line) + 1].strip()
                                        break

                            summary_parts.append(f"  Slide {slide_num} ({title}): Referenced {count} times")

            return summary_parts
